# Rohan Putcha
# required module: python-pptx

import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders

from tkinter import filedialog
from tkinter import *
import os

from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


class Application(Frame):
    def __init__(self, master):
        Frame.__init__(self, master)
        Frame.configure(self, bg="LightBlue1")
        self.grid()
        self.worked = True
        self.prs = Presentation()
        self.prs.save("report.pptx")
        self.createwidgets()

    def createwidgets(self):
        self.firstline = Label(self, bg="LightBlue1")
        self.firstline.grid(row=0, column=0)
        self.label1 = Label(self, text="  Select the master folder containing all required files\t", bg="LightBlue1")
        self.label1.grid(row=1, column=0, sticky=E)
        self.browse = Button(self, text="Browse", command=self.opendir, bg="SkyBlue3", font=("Arial", 9, "bold"))
        self.browse.grid(row=1, column=1, sticky=W)
        self.rightspace = Label(self, text="    ", bg="LightBlue1")
        self.rightspace.grid(row=1, column=2)
        self.chosendirect = Label(self, text="  Selected Directory: ", font=("Arial", 9, "bold"), bg="LightBlue1")
        self.chosendirect.grid(row=2, column=0, sticky=W)
        self.secondline = Label(self, bg="LightBlue1")
        self.secondline.grid(row=3, column=0)
        self.emailaskline = Label(self, text="  Email address to send report: ", font=("Arial", 9, "bold"), bg="LightBlue1")
        self.emailaskline.grid(row=4, column=0, sticky=W)
        self.emailbox = Entry(self, bg="LightBlue1")
        self.emailbox.grid(row=4, column=1, sticky=W)
        self.thirdline = Label(self, bg="LightBlue1")
        self.thirdline.grid(row=5, column=0)
        self.label2 = Label(self, text="Generate Report\t ", bg="LightBlue1")
        self.label2.grid(row=6, column=0, sticky=E)
        self.apply = Button(self, text="Apply", command=self.textfile, bg="SkyBlue3", font=("Arial", 9, "bold"))
        self.apply.grid(row=6, column=1, sticky=W)
        self.fourthline = Label(self, bg="LightBlue1")
        self.fourthline.grid(row=7, column=0)
    def opendir(self):
        self.directory = filedialog.askdirectory(initialdir="/",  title='Select the folder')
        self.chosendirect.configure(text="  Selected Directory: "+self.directory, font=("Arial", 9, "bold"))
        self.chosendirect.grid(row=2, column=0, sticky=W, columnspan=2)

    def readfile(self, direct):
        file = open(direct, "r")
        filelines = file.readlines()
        neededlines = []
        linenum = 0
        for line in filelines:
            for char in filelines[linenum]:
                if char == "G":
                    neededlines.append(line)
            linenum += 1

        grid_id = []
        x = []
        y = []
        z = []
        for i in neededlines:
            grid_id.append(i.split()[0])
            x.append(i.split()[2])
            y.append(i.split()[3])
            z.append(i.split()[4])
        data = [grid_id, x, y, z]
        slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        xi, yi, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)
        shape = slide1.shapes.add_table(5, 4, xi, yi, cx, cy)
        slide1.shapes.title.text = 'Data Table ('+self.filename+')'
        table = shape.table
        table.cell(0, 0).text = "Grid ID"
        table.cell(0, 1).text = "X Displacement"
        table.cell(0, 2).text = "Y Displacement"
        table.cell(0, 3).text = "Z Displacement"

        row = 1
        col = 0
        for val1, val2, val3, val4 in data:
            table.cell(row, col).text = (val1.lower())
            table.cell(row+1, col).text = (val2.lower())
            table.cell(row+2, col).text = (val3.lower())
            table.cell(row+3, col).text = (val4.lower())
            col += 1

        slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        xChart = ChartData()
        xChart.categories = grid_id
        xChart.add_series('X Displacement', tuple(x))
        chart = slide2.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1.7), Inches(0), Inches(8), Inches(2.4), xChart).chart
        chart.series[0].smooth = True

        yChart = ChartData()
        yChart.categories = grid_id
        yChart.add_series('Y Displacement', tuple(y))
        chart = slide2.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1.7), Inches(2.4), Inches(8), Inches(2.4), yChart).chart
        chart.series[0].smooth = True

        zChart = ChartData()
        zChart.categories = grid_id
        zChart.add_series('Z Displacement', tuple(z))
        chart = slide2.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1.7), Inches(4.8), Inches(8), Inches(2.4), zChart).chart
        chart.series[0].smooth = True

    def emailFile(self):
        print("Completed report successfully")
        self.prs.save(self.directory+"/report.pptx")
        email = self.emailbox.get()
        email_user = 'fcaprojectmail@gmail.com'
        email_password = 'projectreport'
        email_send = email
        subject = 'Automated Project Report'
        message = MIMEMultipart()
        body = 'The project report is attached.'

        message['From'] = email_user
        message['To'] = email_send
        message['Subject'] = subject

        message.attach(MIMEText(body, 'plain'))
        filename = self.directory+'/report.pptx'
        attachment = open(filename, 'rb')
        mimebase = MIMEBase('application', 'octet-stream')
        mimebase.set_payload(attachment.read())
        encoders.encode_base64(mimebase)
        mimebase.add_header('Content-Disposition', "attachment; filename= " + 'report.ppt')
        message.attach(mimebase)
        text = message.as_string()

        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()
        server.login(email_user, email_password)
        server.sendmail(email_user, email_send, text)
        server.quit()

        print('Emailed report to ' + email + ' successfully.')
    def textfile(self):
        try:
            a = os.listdir(self.directory)
            for subfolder in a:
                directstring = self.directory
                directstring += "/" + subfolder
                temp = directstring
                b = os.listdir(directstring)
                for file in b:
                    directstring = temp
                    directstring += "/" + file
                    self.filename = file
                    print("Opening directory \""+directstring+"\"...")
                    self.readfile(directstring)
            self.worked = True
        except NotADirectoryError:
            print("An error occurred. Make sure only subfolders are in the master folder.")
            self.worked = False
        except:
            print("An error occurred. Make sure you've provided the correct path.")
            self.worked = False
        if self.worked:
            self.emailFile()

root = Tk()
root.title("Report Generation Tool")
root.configure(bg="LightBlue1")
app = Application(root)
root.mainloop()
